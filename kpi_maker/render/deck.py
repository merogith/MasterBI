"""Board deck (PPTX).

The one rule that makes a deck read like a consulting deliverable: **the slide
title is the message, not the topic.** "Churn is concentrated in SMB" — never
"Churn analysis". Every exhibit slide here takes its headline from the finding
attached to that chart, so the titles say something by construction.
"""
from __future__ import annotations

import io
from pathlib import Path
from typing import Dict, List, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

from ..fmt import fmt_value
from ..insight.detectors import Finding
from ..kpi.schema import KPISet
from ..metrics.engine import MetricResult
from ..profile.schema import CompanyProfile
from ..viz.theme import STATUS_LABEL, TOKENS
from .sections import SectionContext
from .sections import build as build_sections

# The default only. `brand.font_stack` overrides it per run — PPTX stores a
# font name and lets PowerPoint resolve it, so there is no file to find.
DEFAULT_FONT = "Segoe UI"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.6)


def _rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


class Deck:
    def __init__(self, profile: CompanyProfile,
                 tokens: Optional[Dict[str, str]] = None,
                 font: Optional[str] = None,
                 footer_text: Optional[str] = None,
                 locale: Optional[str] = None):
        self.t = dict(tokens or TOKENS["light"])
        self.font = font or DEFAULT_FONT
        self.footer_text = footer_text
        self.locale = locale
        self.logo = None
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_W
        self.prs.slide_height = SLIDE_H
        self.profile = profile
        self.blank = self.prs.slide_layouts[6]

    def _slide(self):
        slide = self.prs.slides.add_slide(self.blank)
        self.last = slide
        return slide

    def speaker_note(self, paragraphs: List[str]) -> None:
        """Narrative prose belongs in the notes pane, not on the slide.

        This is the one place the deck treats the AI paragraph differently from
        the PDF and the DOCX, and it is not a compromise. Connective prose is
        what a presenter *says* while a slide is up — putting it on the slide
        would mean either shrinking it to a bullet, which loses the connective
        work, or a wall of text on a deck whose whole discipline is one message
        per slide. The notes pane is where that sentence already lives.
        """
        slide = getattr(self, "last", None)
        if slide is None or not paragraphs:
            return
        slide.notes_slide.notes_text_frame.text = "\n\n".join(paragraphs)

    def _text(self, slide, text, left, top, width, height, size=14, bold=False,
              color="text_primary", align=PP_ALIGN.LEFT, wrap=True):
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = wrap
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = self.font
        run.font.color.rgb = _rgb(self.t[color])
        return box

    def _headline(self, slide, message: str, kicker: str = ""):
        """The message, not the topic."""
        if kicker:
            self._text(slide, kicker.upper(), MARGIN, Inches(0.42),
                       SLIDE_W - 2 * MARGIN, Inches(0.28), size=10, color="muted")
        self._text(slide, message, MARGIN, Inches(0.72),
                   SLIDE_W - 2 * MARGIN, Inches(0.9), size=23, bold=True)
        line = slide.shapes.add_shape(1, MARGIN, Inches(1.62),
                                      SLIDE_W - 2 * MARGIN, Emu(9525))
        line.fill.solid()
        line.fill.fore_color.rgb = _rgb(self.t["grid"])
        line.line.fill.background()
        line.shadow.inherit = False

    def _footer(self, slide, text: str):
        self._text(slide, text, MARGIN, SLIDE_H - Inches(0.52),
                   SLIDE_W - 2 * MARGIN, Inches(0.3), size=9, color="muted")

    # -- slide types ------------------------------------------------------
    def title_slide(self, results: List[MetricResult], kpi_set: KPISet, period: str):
        s = self._slide()
        p = self.profile
        if self.logo is not None:
            s.shapes.add_picture(io.BytesIO(self.logo.data), MARGIN,
                                 Inches(1.2), height=Inches(0.55))
        self._text(s, "PERFORMANCE REVIEW", MARGIN, Inches(2.3),
                   Inches(8), Inches(0.3), size=11, color="muted")
        self._text(s, p.identity.name, MARGIN, Inches(2.8),
                   Inches(11), Inches(1.0), size=40, bold=True)
        self._text(s, f"{p.business_model.type.value.upper()} · "
                      f"{p.business_model.customer_type.value} · "
                      f"{p.identity.country} · {period}",
                   MARGIN, Inches(3.9), Inches(11), Inches(0.4), size=14,
                   color="text_secondary")
        north = next((r for r in results if r.kpi.id == kpi_set.north_star and r.computed), None)
        if north:
            self._text(s, north.kpi.name.upper(), MARGIN, Inches(4.9),
                       Inches(8), Inches(0.3), size=10, color="muted")
            self._text(s, fmt_value(north.current, north.kpi.unit,
                                    p.identity.currency, locale=self.locale),
                       MARGIN, Inches(5.2), Inches(8), Inches(0.9), size=34,
                       bold=True, color="series_1")
        self._footer(s, self.footer_text or
                     f"Prepared for the {p.intent.audience.value} · "
                     f"benchmarks are illustrative — see appendix")

    def bullets_slide(self, message: str, items: List[str], kicker: str = ""):
        s = self._slide()
        self._headline(s, message, kicker)
        top = Inches(2.0)
        box = s.shapes.add_textbox(MARGIN, top, SLIDE_W - 2 * MARGIN,
                                   SLIDE_H - top - Inches(0.8))
        tf = box.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            run = p.add_run()
            run.text = item
            run.font.size = Pt(14)
            run.font.name = self.font
            run.font.color.rgb = _rgb(self.t["text_secondary"])
            p.space_after = Pt(11)
        return s

    def exhibit_slide(self, message: str, png: Optional[bytes], kicker: str = "",
                      caption: str = ""):
        if png is None:
            return None
        s = self._slide()
        self._headline(s, message, kicker)
        img_w = SLIDE_W - 2 * MARGIN
        # Charts export at 900x360 or 900x320; scale to width and centre.
        pic = s.shapes.add_picture(io.BytesIO(png), MARGIN, Inches(2.0), width=img_w)
        max_h = SLIDE_H - Inches(2.0) - Inches(0.9)
        if pic.height > max_h:
            ratio = max_h / pic.height
            pic.height = int(max_h)
            pic.width = int(pic.width * ratio)
            pic.left = int((SLIDE_W - pic.width) / 2)
        if caption:
            self._footer(s, caption)
        return s

    def table_slide(self, message: str, headers: List[str], rows: List[List[str]],
                    kicker: str = "", col_widths: Optional[List[float]] = None):
        s = self._slide()
        self._headline(s, message, kicker)
        rows = rows[:12]
        n_rows, n_cols = len(rows) + 1, len(headers)
        top = Inches(2.0)
        height = min(SLIDE_H - top - Inches(0.7), Inches(0.32) * n_rows)
        shape = s.shapes.add_table(n_rows, n_cols, MARGIN, top,
                                   SLIDE_W - 2 * MARGIN, height)
        table = shape.table
        if col_widths:
            total = SLIDE_W - 2 * MARGIN
            for i, frac in enumerate(col_widths):
                table.columns[i].width = Emu(int(total * frac))

        for c, head in enumerate(headers):
            cell = table.cell(0, c)
            cell.text = head
            para = cell.text_frame.paragraphs[0]
            para.runs[0].font.size = Pt(10)
            para.runs[0].font.bold = True
            para.runs[0].font.name = self.font
            para.runs[0].font.color.rgb = _rgb(self.t["text_primary"])

        for r, row in enumerate(rows, start=1):
            for c, val in enumerate(row):
                cell = table.cell(r, c)
                cell.text = str(val)
                para = cell.text_frame.paragraphs[0]
                para.runs[0].font.size = Pt(10)
                para.runs[0].font.name = self.font
                para.runs[0].font.color.rgb = _rgb(
                    self.t["text_primary"] if c == 0 else self.t["text_secondary"])
        return s

    def save(self, path: Path) -> Path:
        path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(path))
        return path


# The deck shares the report's *content* for the three sections where the two
# genuinely say the same thing — the executive summary, the risks and the
# actions. It keeps its own title slide, its own top-tier scorecard and its own
# exhibit plan, because those are not the report's sections rendered smaller:
# the scorecard shows only tier-1 KPIs with five columns, and every exhibit
# gets a headline written from its finding rather than the chart's topic. That
# is a different document, and forcing it through the shared registry would
# either change the deck or fill the registry with per-format branches.
#
# Section toggles still reach all of it: a section switched off in the spec
# produces no slide here either.

# Fewer rows than a page holds. Recorded here rather than left as a bare 8.
DECK_LIMITS = {"exec_summary": {"limit": 5}, "risks": {"limit": 5},
               "actions": {"limit": 8}}

# Which section owns each planned exhibit, so that disabling a section drops
# its charts from the deck too.
EXHIBIT_PLAN = [
    ("arr_trend", "arr", "Trajectory", "deep_dives"),
    ("arr_bridge", "net_new_arr", "Diagnostic", "diagnostic"),
    ("retention", "nrr", "Retention", "deep_dives"),
    ("segment_churn", "logo_churn_rate", "Retention", "deep_dives"),
    ("cohort_heatmap", "grr", "Retention", "diagnostic"),
    ("cac_payback", "cac_payback_months", "Efficiency", "deep_dives"),
    ("channel_cost", "blended_cac", "Efficiency", "deep_dives"),
    ("indexed_growth", "arr_per_fte", "Leverage", "deep_dives"),
    ("benchmark_position", None, "Benchmarks", "benchmarks"),
]
EXHIBIT_SECTIONS = ("diagnostic", "deep_dives", "benchmarks")

#: Named here for the callers that already say `STATUS_WORD`; the map
#: itself is `viz.theme`'s, so the three renderers cannot drift apart.
STATUS_WORD = STATUS_LABEL


def _scorecard_slide(deck: Deck, results: List[MetricResult], cur: str) -> None:
    computed = [r for r in results if r.computed]
    top = [r for r in computed if int(r.kpi.tier) <= 1]
    n_red = len([r for r in computed if r.status == "red"])
    n_amber = len([r for r in computed if r.status == "amber"])
    deck.table_slide(
        f"{n_red} KPI{'' if n_red == 1 else 's'} off track, {n_amber} on watch",
        ["KPI", "Current", "12mo ago", "Target", "Status"],
        [[r.kpi.short_name or r.kpi.name,
          fmt_value(r.current, r.kpi.unit, cur, locale=deck.locale),
          fmt_value(r.prior_year, r.kpi.unit, cur, locale=deck.locale),
          fmt_value(r.target, r.kpi.unit, cur, locale=deck.locale),
          STATUS_WORD.get(r.status, "—")] for r in top],
        kicker="Scorecard",
        col_widths=[0.36, 0.16, 0.16, 0.16, 0.16],
    )


def _exhibit_slides(deck: Deck, kpi_set: KPISet, results: List[MetricResult],
                    findings: List[Finding], specs: List,
                    images: Dict[str, bytes], enabled: set, cur: str) -> None:
    """One slide per chart, headlined by the finding attached to it."""
    finding_for: Dict[str, Finding] = {}
    for f in findings:
        for kid in f.kpi_ids:
            finding_for.setdefault(kid, f)

    # Where no finding is attached, compute a headline rather than falling back
    # to the chart's topic. "Annual Recurring Revenue" tells a board nothing.
    computed = {r.kpi.id: r for r in results if r.computed}
    fallbacks: Dict[str, str] = {}
    north, growth = computed.get(kpi_set.north_star), computed.get("arr_growth_yoy")
    if north and north.current is not None:
        headline = (f"{north.kpi.name} reached "
                    f"{fmt_value(north.current, north.kpi.unit, cur, locale=deck.locale)}")
        if growth and growth.current is not None:
            headline += f", up {growth.current:.0%} year on year"
        fallbacks["arr_trend"] = headline
    benched = [r for r in computed.values() if r.benchmark_position]
    # `outside_band` belongs here: a target_band metric outside its cohort's
    # middle half is behind the cohort, even though neither "below median" nor
    # "bottom quartile" can be said of a metric with no better end.
    behind = [r for r in benched
              if r.benchmark_position in ("below_median", "bottom_quartile",
                                          "outside_band")]
    if benched:
        fallbacks["benchmark_position"] = (
            f"{len(behind)} of {len(benched)} benchmarked KPIs sit behind the "
            f"cohort")

    specs_by_id = {s.id: s for s in specs}
    for spec_id, kpi_id, kicker, owner in EXHIBIT_PLAN:
        if owner not in enabled:
            continue
        spec = specs_by_id.get(spec_id)
        if spec is None or spec_id not in images:
            continue
        f = finding_for.get(kpi_id) if kpi_id else None
        headline = f.title if f else fallbacks.get(spec_id, spec.title)
        deck.exhibit_slide(headline, images[spec_id], kicker=kicker,
                           caption=spec.note or spec.subtitle)


def render_deck(path: Path, profile: CompanyProfile, kpi_set: KPISet,
                results: List[MetricResult], findings: List[Finding],
                images: Dict[str, bytes], specs: List, period: str = "",
                tokens: Optional[Dict[str, str]] = None,
                section_order: Optional[List[str]] = None,
                logo=None,
                narrative: Optional[Dict[str, List[str]]] = None,
                caveats: Optional[List[str]] = None,
                font_stack: Optional[List[str]] = None,
                footer_text: Optional[str] = None,
                locale: Optional[str] = None) -> Path:
    deck = Deck(profile, tokens, font=(font_stack or [None])[0],
                footer_text=footer_text, locale=locale)
    deck.logo = logo
    cur = profile.identity.currency

    ctx = SectionContext(
        profile=profile, kpi_set=kpi_set, results=results, findings=findings,
        images=images, specs=specs, period=period,
        narrated=sorted(narrative or {}), caveats=list(caveats or []),
        locale=locale)
    contents = build_sections(ctx, section_order, limits=DECK_LIMITS,
                              narrative=narrative)
    enabled = {c.id for c in contents}
    by_id = {c.id: c for c in contents}
    exhibits_done = False

    for content in contents:
        if content.id == "cover":
            deck.title_slide(results, kpi_set, period)

        elif content.id == "exec_summary":
            deck.bullets_slide(
                "What the numbers say",
                [f"{b.lead}\n{b.text}" for b in content.bullets],
                kicker="Executive summary")

        elif content.id == "scorecard":
            _scorecard_slide(deck, results, cur)

        elif content.id in EXHIBIT_SECTIONS:
            # The plan spans three sections, so it runs once, at the first of
            # them that survived the spec.
            if not exhibits_done:
                exhibits_done = True
                _exhibit_slides(deck, kpi_set, results, findings, specs, images,
                                enabled, cur)

        elif content.id == "risks":
            if content.bullets:
                deck.bullets_slide(
                    f"{content.total} issues need a decision this quarter",
                    # No severity prefix: the headline already says these are
                    # the risks, and the slide has less room to spend on it.
                    [f"{b.title}\n{b.text}" for b in content.bullets],
                    kicker="Risks and watch-list")

        elif content.id == "actions":
            for table in content.tables:
                deck.table_slide(
                    "Where to act first",
                    ["Action", "Owner", "Impact", "Effort"],
                    # The deck drops the "Moves" column the page has room for.
                    [[row[0][:110], row[2], row[3], row[4]] for row in table.rows],
                    kicker="Recommended actions",
                    col_widths=[0.58, 0.14, 0.14, 0.14])

        # Attached to whichever slide this section finished on, which is the
        # one a presenter will be looking at when they need the sentence.
        deck.speaker_note(content.narrative)

    appendix = by_id.get("appendix")
    if appendix is not None:
        deck.bullets_slide(
            "How to read this deck",
            [
                "Every figure is computed by deterministic code from the underlying "
                "fact tables. No number in this deck was produced by a language model.",
                f"{len(kpi_set.kpis)} KPIs were selected from the "
                f"{profile.business_model.type.value} library by scoring applicability, "
                f"objective alignment, Balanced Scorecard coverage and audience fit.",
                "Benchmarks are illustrative placeholders assembled from public "
                "commentary — not a licensed dataset. They are suitable for "
                "calibration, not for external reporting.",
                f"Profile confidence {profile.confidence:.0%}. Fields filled from "
                f"sector defaults rather than your data are listed in the appendix "
                f"of the full report.",
            ],
            kicker="Appendix")

    return deck.save(path)
